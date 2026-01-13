from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Cores
    BG_COLOR = RGBColor(19, 23, 34)    # Azul escuro quase preto
    TEXT_MAIN = RGBColor(255, 255, 255) # Branco
    ACCENT = RGBColor(79, 70, 229)      # Indigo vibrante
    SUBTEXT = RGBColor(203, 213, 225)   # Cinza claro

    # Função auxiliar para configurar fundo do slide
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # --- SLIDE 1: TÍTULO ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blank
    set_dark_background(slide)
    
    # Barra lateral decorativa
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # Título Grande
    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(8), Inches(2))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "CAFÉ POINT"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = 'Arial'

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(8), Inches(1))
    stf = sub_box.text_frame
    sp = stf.add_paragraph()
    sp.text = "O Futuro da Gestão de Restaurantes"
    sp.font.size = Pt(28)
    sp.font.color.rgb = TEXT_MAIN
    sp.font.name = 'Arial'
    
    # Rodapé
    foot_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(5), Inches(1))
    ftf = foot_box.text_frame
    fp = ftf.add_paragraph()
    fp.text = "Apresentação Executiva 2026"
    fp.font.size = Pt(14)
    fp.font.color.rgb = SUBTEXT

    # --- FUNÇÃO PARA SLIDES DE CONTEÚDO ---
    def add_content_slide(title, main_point, bullet_points=[]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_dark_background(slide)
        
        # Barra superior colorida
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(1.5), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Título
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
        tf = tbox.text_frame
        p = tf.add_paragraph()
        p.text = title.upper()
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.font.name = 'Arial'

        # Ponto Principal (Destaque)
        mbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        mf = mbox.text_frame
        mf.word_wrap = True
        mp = mf.add_paragraph()
        mp.text = main_point
        mp.font.size = Pt(24)
        mp.font.color.rgb = SUBTEXT
        mp.font.name = 'Arial'
        
        # Bullets (se houver)
        if bullet_points:
            bbox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8.5), Inches(3))
            bf = bbox.text_frame
            bf.word_wrap = True
            for point in bullet_points:
                bp = bf.add_paragraph()
                bp.text = "• " + point
                bp.font.size = Pt(20)
                bp.font.color.rgb = TEXT_MAIN
                bp.space_after = Pt(14)

    # --- CRIAÇÃO DOS SLIDES ---

    add_content_slide(
        "O Problema",
        "A ineficiência operacional está custando dinheiro e clientes todos os dias.",
        [
            "Pedidos no papel causam erros e atrasos na cozinha.",
            "Falta de visibilidade do stock em tempo real.",
            "Dificuldade em medir lucros e perdas por prato."
        ]
    )

    add_content_slide(
        "A Solução",
        "Uma plataforma digital unificada que conecta todos os pontos do restaurante.",
        [
            "POS Digital no Salão (Tablet/PC).",
            "KDS (Ecrã de Cozinha) automatizado.",
            "Backoffice financeiro integrado."
        ]
    )

    add_content_slide(
        "Eficiência no Salão",
        "Aumente a rotação de mesas em até 30% com pedidos instantâneos.",
        [
            "Status de mesas em tempo real (Sem gritos).",
            "Envio de pedidos directo para as estações de preparo.",
            "Redução drástica de erros de anotação."
        ]
    )
    
    add_content_slide(
        "Cozinha Inteligente",
        "Organize o caos e garanta que os pratos saiam na ordem certa.",
        [
            "Priorização automática de pedidos.",
            "Métricas de tempo de preparo.",
            "Comunicação visual clara e silenciosa."
        ]
    )
    
    add_content_slide(
        "Controlo Total",
        "Transforme stock em dinheiro e pare de adivinhar os custos.",
        [
            "Baixa automática de ingredientes por ficha técnica.",
            "Alertas de stock mínimo para compras.",
            "Relatórios de vendas detalhados por dia/mês."
        ]
    )
    
    add_content_slide(
        "Próximos Passos",
        "Plano de implementação rápida para resultados imediatos.",
        [
            "Semana 1: Configuração e Piloto.",
            "Semana 2: Formação da Equipa.",
            "Semana 3: Lançamento Oficial (Go-Live)."
        ]
    )

    # --- SLIDE FINAL ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    end_box = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(10), Inches(2))
    ef = end_box.text_frame
    ep = ef.add_paragraph()
    ep.text = "Obrigado"
    ep.font.size = Pt(54)
    ep.font.bold = True
    ep.font.color.rgb = ACCENT
    ep.alignment = PP_ALIGN.CENTER
    
    ep2 = ef.add_paragraph()
    ep2.text = "cafepoint.sistema"
    ep2.font.size = Pt(20)
    ep2.font.color.rgb = SUBTEXT
    ep2.alignment = PP_ALIGN.CENTER

    prs.save("Cafe_Point_Presentation_V2_Dark.pptx")
    print("Nova apresentação (visual Dark Mode) gerada com sucesso!")

if __name__ == "__main__":
    create_presentation()
