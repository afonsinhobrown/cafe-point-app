from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation_v3():
    prs = Presentation()
    
    # --- COLORS DEFINED FROM apresentacao.html ---
    PRIMARY_INDIGO = RGBColor(79, 70, 229)   # #4f46e5
    SECONDARY_DARK = RGBColor(30, 27, 75)    # #1e1b4b
    ACCENT_AMBER = RGBColor(245, 158, 11)    # #f59e0b
    TEXT_DARK = RGBColor(51, 51, 51)         # #333333
    BG_LIGHT = RGBColor(248, 250, 252)       # #f8fafc
    WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # --- SLIDE 1: TITLE (Mimicking .intro class) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, SECONDARY_DARK)
    
    # Decorative Accent Line (mimicking border-bottom)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(10), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_INDIGO
    shape.line.fill.background()

    # Title
    tbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = tbox.text_frame
    p = tf.add_paragraph()
    p.text = "CAFÉ POINT"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Subtitle
    sbox = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    sf = sbox.text_frame
    sp = sf.add_paragraph()
    sp.text = "Sistema Integrado de Gestão"
    sp.font.size = Pt(32)
    sp.font.color.rgb = ACCENT_AMBER
    sp.alignment = PP_ALIGN.CENTER
    sp.font.name = "Segoe UI"

    # Tagline
    tagbox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tagf = tagbox.text_frame
    tagp = tagf.add_paragraph()
    tagp.text = "Modernização, Controlo e Eficiência Operacional"
    tagp.font.size = Pt(20)
    tagp.font.color.rgb = RGBColor(226, 232, 240) # Light gray
    tagp.alignment = PP_ALIGN.CENTER
    tagp.font.name = "Segoe UI"

    # --- HELPER FOR CONTENT SLIDES ---
    def add_content_slide(title, points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, BG_LIGHT)

        # Header Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_INDIGO
        bar.line.fill.background()

        # Title
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tbox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_DARK
        p.font.name = "Segoe UI"
        
        # Underline accent
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_AMBER
        line.line.fill.background()

        # Content Box
        cbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
        cf = cbox.text_frame
        cf.word_wrap = True

        for point in points:
            # Bullet point symbol (simulated)
            cp = cf.add_paragraph()
            cp.text = f"{point}" 
            cp.font.size = Pt(22)
            cp.font.color.rgb = TEXT_DARK
            cp.font.name = "Segoe UI"
            cp.space_after = Pt(20)
            cp.level = 0
            
            # Simple manual bullet simulation if needed, but standard Paragraph level usually handles it.
            # To make it look like the prompt's rich text, we rely on the text being descriptive.

    # --- SLIDES GENERATION ---

    # Slide 2: Problemas
    add_content_slide("Desafios Operacionais Atuais", [
        "🛑 Ineficiência: Pedidos em papel causam erros e atrasos na comunicação.",
        "🛑 Quebras de Stock: Falta de rastreabilidade gera desperdícios.",
        "🛑 Falta de Dados: Gestão baseada em 'feeling', sem relatórios precisos.",
        "🛑 Experiência do Cliente: Tempo de espera elevado afecta a satisfação."
    ])

    # Slide 3: Solução
    add_content_slide("Solução: Café Point", [
        "Uma plataforma 'All-in-One' que conecta Salão, Cozinha e Backoffice.",
        "🎯 Foco: Eliminar papel, automatizar processos e garantir controlo.",
        "💻 Tecnologia: Sistema moderno, seguro e acessível via Tablets/PC."
    ])

    # Slide 4: POS
    add_content_slide("Atendimento Ágil (POS)", [
        "✅ Mapa de Mesas Digital: Visualização em tempo real (Livre/Ocupada).",
        "✅ Pedido Mobile: Garçom lança o pedido na mesa.",
        "✅ Personalização: Adição fácil de observações (ex: 'sem gelo')."
    ])

    # Slide 5: KDS
    add_content_slide("Cozinha Conectada (KDS)", [
        "👨‍🍳 Fim das 'Bonitas': Pedidos aparecem no ecrã da cozinha.",
        "⏱️ Controlo de Tempo: Cozinheiros sabem exactamante o que preparar.",
        "🔔 Status: Fluxo claro de Pendente -> Preparando -> Pronto."
    ])
    
    # Slide 6: Stock
    add_content_slide("Gestão de Stock", [
        "📦 Ficha Técnica: Baixa automática de ingredientes ao vender.",
        "⚠️ Alertas Inteligentes: Aviso automático de stock mínimo.",
        "📊 Histórico: Rastreio completo de todas as entradas e saídas."
    ])

    # Slide 7: ROI
    add_content_slide("Impacto Esperado (ROI)", [
        "🚀 Aumento de 30% na rotação de mesas.",
        "💰 Redução de 15% em desperdícios.",
        "⭐ Melhor experiência do cliente (menos erros).",
        "📈 Decisões baseadas em dados reais."
    ])

    # Slide 8: Next Steps
    add_content_slide("Roteiro de Implementação", [
        "1. Semana 1: Instalação Piloto e Configuração.",
        "2. Semana 2: Treinamento da Equipe.",
        "3. Semana 3: 'Go-Live' assistido.",
        "👉 Aprovação para iniciar o piloto."
    ])

    # --- FINAL SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, SECONDARY_DARK)
    
    # Center Box
    fbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    ff = fbox.text_frame
    fp = ff.add_paragraph()
    fp.text = "Obrigado!"
    fp.font.size = Pt(54)
    fp.font.bold = True
    fp.font.color.rgb = WHITE
    fp.alignment = PP_ALIGN.CENTER
    
    fp2 = ff.add_paragraph()
    fp2.text = "Café Point 2026"
    fp2.font.size = Pt(24)
    fp2.font.color.rgb = ACCENT_AMBER
    fp2.alignment = PP_ALIGN.CENTER
    
    # Save
    prs.save("Cafe_Point_Presentation_V3_Branded.pptx")
    print("Apresentação V3 (Branded) gerada com sucesso!")

if __name__ == "__main__":
    create_presentation_v3()
