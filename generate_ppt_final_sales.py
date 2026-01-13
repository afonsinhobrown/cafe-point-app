from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sales_presentation():
    prs = Presentation()

    # --- PALETA DE CORES (PREMIUM SAAS) ---
    BG_DARK = RGBColor(15, 23, 42)      # #0f172a (Slate 900)
    TEXT_WHITE = RGBColor(248, 250, 252) # Slate 50
    TEXT_GRAY = RGBColor(148, 163, 184)  # Slate 400
    ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo 500
    ACCENT_RED = RGBColor(239, 68, 68)   # Red 500 (Para problemas)
    ACCENT_GREEN = RGBColor(34, 197, 94) # Green 500 (Para lucro)

    def set_slide_style(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    # --- SLIDE 1: CAPA IMPACTANTE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide)

    # Título Grande
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "CAFÉ POINT"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = 'Arial'
    
    # Subtítulo (Proposta de Valor)
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
    sf = sub_box.text_frame
    sp = sf.add_paragraph()
    sp.text = "O Fim do Caos Operacional no Seu Restaurante."
    sp.font.size = Pt(28)
    sp.font.color.rgb = ACCENT_INDIGO
    sp.font.name = 'Arial'

    # --- SLIDE 2: O PROBLEMA (DOR) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide)

    # Título
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tbox.text_frame
    p = tf.add_paragraph()
    p.text = "QUANTO DINHEIRO ESTAMOS PERDENDO?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.font.name = 'Arial'

    # Conteúdo (Grid visual simples)
    cbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
    cf = cbox.text_frame
    cf.word_wrap = True

    points = [
        ("DESPERDÍCIO DE STOCK (20%)", "Sem controlo rigoroso, ingredientes somem ou estragam. O lucro vai para o lixo antes de chegar ao prato."),
        ("ERROS DE PEDIDO", "Letra ilegível e falhas de comunicação cozinha-garçom geram pratos devolvidos e clientes insatisfeitos."),
        ("LENTIDÃO NO ATENDIMENTO", "Cada minuto de atraso é uma mesa que roda menos vezes na noite.")
    ]

    for title, desc in points:
        tp = cf.add_paragraph()
        tp.text = title
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = TEXT_WHITE
        tp.space_before = Pt(20)

        dp = cf.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(18)
        dp.font.color.rgb = TEXT_GRAY
        dp.space_after = Pt(10)

    # --- SLIDE 3: A SOLUÇÃO (HERÓI) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide)

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tbox.text_frame
    p = tf.add_paragraph()
    p.text = "A SOLUÇÃO: CONTROLO TOTAL"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.font.name = 'Arial'

    sbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    sf = sbox.text_frame
    sp = sf.add_paragraph()
    sp.text = "Uma plataforma única que conecta tudo em tempo real."
    sp.font.size = Pt(24)
    sp.font.color.rgb = TEXT_WHITE

    # Funcionalidades (Boxed style simulado)
    fbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4))
    ff = fbox.text_frame
    
    feats = [
        ("📱 POS MÓVEL", "O garçom lança o pedido na mesa. Zero erros. Zero deslocações inúteis."),
        ("👨‍🍳 KDS (COZINHA DIGITAL)", "Ecrãs substituem papel. Fila organizada por ordem de chegada."),
        ("📉 STOCK AUTOMÁTICO", "Vendeu um prato? O sistema baixa os ingredientes. Instantâneo.")
    ]

    for title, desc in feats:
        tp = ff.add_paragraph()
        tp.text = "• " + title
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = ACCENT_INDIGO
        tp.space_before = Pt(24)
        
        dp = ff.add_paragraph()
        dp.text = "   " + desc
        dp.font.size = Pt(20)
        dp.font.color.rgb = TEXT_GRAY

    # --- SLIDE 4: ROI (DINHEIRO) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide)

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tbox.text_frame
    p = tf.add_paragraph()
    p.text = "O IMPACTO FINANCEIRO (ROI)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = 'Arial'

    cbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    cf = cbox.text_frame
    
    # Coluna 1: Aumentar Receita
    p1 = cf.add_paragraph()
    p1.text = "🚀 AUMENTO DE RECEITA"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(10)

    items1 = ["+ Rotação de Mesas", "+ Ticket Médio (Upsell)", "+ Fidelização de Clientes"]
    for i in items1:
        ip = cf.add_paragraph()
        ip.text = "• " + i
        ip.font.size = Pt(20)
        ip.font.color.rgb = TEXT_GRAY

    # Coluna 2: Reduzir Custos
    cbox2 = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
    cf2 = cbox2.text_frame
    
    p2 = cf2.add_paragraph()
    p2.text = "💰 REDUÇÃO DE CUSTOS"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(10)

    items2 = ["- Desperdício de Alimentos", "- Roubos e Desvios", "- Erros Operacionais"]
    for i in items2:
        ip = cf2.add_paragraph()
        ip.text = "• " + i
        ip.font.size = Pt(20)
        ip.font.color.rgb = TEXT_GRAY

    # --- SLIDE 5: CHAMADA PARA AÇÃO ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide)

    end_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    ef = end_box.text_frame
    ep = ef.add_paragraph()
    ep.text = "VAMOS MODERNIZAR?"
    ep.font.size = Pt(48)
    ep.font.bold = True
    ep.font.color.rgb = ACCENT_INDIGO
    ep.alignment = PP_ALIGN.CENTER
    
    ep2 = ef.add_paragraph()
    ep2.text = "Implementação completa em 3 semanas."
    ep2.font.size = Pt(24)
    ep2.font.color.rgb = TEXT_WHITE
    ep2.alignment = PP_ALIGN.CENTER
    ep2.space_before = Pt(20)

    ep3 = ef.add_paragraph()
    ep3.text = "Agende o 'Go-Live' hoje."
    ep3.font.size = Pt(20)
    ep3.font.color.rgb = TEXT_GRAY
    ep3.alignment = PP_ALIGN.CENTER
    ep3.space_before = Pt(10)

    # Save
    prs.save("Apresentacao_Venda_Cafe_Point.pptx")
    print("PPT de Vendas (ROI Focado) gerado com sucesso!")

if __name__ == "__main__":
    create_sales_presentation()
