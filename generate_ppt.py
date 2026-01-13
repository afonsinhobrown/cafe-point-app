from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Paleta minimalista Apple/SaaS premium
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(60, 60, 60)

def slide_title_only(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER

def slide_statement(title, statement):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(7.5), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BLACK

    sbox = slide.shapes.add_textbox(Inches(1.2), Inches(3), Inches(7.5), Inches(2))
    stf = sbox.text_frame
    sp = stf.paragraphs[0]
    sp.text = statement
    sp.font.size = Pt(24)
    sp.font.color.rgb = DARK_GRAY

# Criar slides
slide_title_only("Café Point", "Gestão moderna para restaurantes modernos")
slide_statement("O Problema", "A maioria dos restaurantes ainda opera com sistemas fragmentados, lentos e sem visibilidade real.")
slide_statement("A Solução", "Uma plataforma única que liga salão, cozinha e gestão em tempo real.")
slide_statement("Atendimento", "Pedidos lançados na mesa. Menos passos. Menos erros. Mais mesas atendidas.")
slide_statement("Cozinha", "Pedidos claros, prioridades visíveis e controlo total do tempo de preparação.")
slide_statement("Stock", "Cada venda actualiza automaticamente custos e inventário.")
slide_statement("Impacto", "Mais eficiência operacional. Menos desperdício. Decisões baseadas em dados.")
slide_statement("Implementação", "Três semanas para entrar em produção com acompanhamento total.")

# Salvar arquivo
prs.save("Cafe_Point_Apple_SaaS_Premium.pptx")
print("PPT gerado com sucesso!")
